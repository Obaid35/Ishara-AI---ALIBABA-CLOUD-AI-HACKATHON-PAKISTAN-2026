"""
Ishara AI - Hackathon Presentation Generator v3
Flawless spacing, no overlapping text, perfect logo contrast & card layouts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
BRAIN = r"C:\Users\Obaid\.gemini\antigravity-ide\brain\e6dcce9a-8ea2-4961-99e0-1d5e973bc78e"
PROJECT = r"d:\Dev\hackathon\psl-bridge-project"
LOGO_LOCKUP = os.path.join(PROJECT, "assets", "logo", "ishara-lockup.png")
LOGO_MARK   = os.path.join(PROJECT, "assets", "logo", "ishara-mark.png")
SS_APP      = os.path.join(BRAIN, "ss_patient.png")
OUTPUT      = os.path.join(PROJECT, "Ishara_AI_Hackathon.pptx")

# ── Colors ──
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT    = RGBColor(0xF8, 0xFA, 0xFC)   # ultra-clean slate-50
GREEN_DARK  = RGBColor(0x01, 0x54, 0x27)
GREEN       = RGBColor(0x01, 0x7A, 0x3A)   # brand green #017A3A
GREEN_TINT  = RGBColor(0xEC, 0xFD, 0xF5)   # light emerald tint
GREEN_BORDER= RGBColor(0xA7, 0xF3, 0xD0)
DARK_NAVY   = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
TEXT_BODY   = RGBColor(0x33, 0x41, 0x55)   # slate-700
TEXT_MUTED  = RGBColor(0x64, 0x74, 0x8B)   # slate-500
RED_BG      = RGBColor(0xFE, 0xF2, 0xF2)
RED_BORDER  = RGBColor(0xFE, 0xCD, 0xCD)
RED_TEXT    = RGBColor(0x99, 0x1B, 0x1B)
AMBER_BG    = RGBColor(0xFF, 0xFB, 0xEB)
AMBER_BORDER= RGBColor(0xFE, 0xE3, 0xA1)
AMBER_DARK  = RGBColor(0xB4, 0x53, 0x09)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def set_bg(slide, color=BG_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, step_num, step_title, main_headline):
    """Clean standard slide header with section pill and headline"""
    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = GREEN; top_bar.line.fill.background()

    # Section pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(0.55), Inches(0.32))
    pill.fill.solid(); pill.fill.fore_color.rgb = GREEN; pill.line.fill.background()
    tf = pill.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = step_num; r.font.size = Pt(12); r.font.color.rgb = WHITE; r.font.bold = True

    # Section title
    tb = slide.shapes.add_textbox(Inches(1.45), Inches(0.42), Inches(6.0), Inches(0.32))
    tf2 = tb.text_frame; tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = step_title; r2.font.size = Pt(12); r2.font.color.rgb = GREEN
    r2.font.bold = True; r2.font.name = "Segoe UI"

    # Main Headline (Single line or clean 24-26pt to NEVER overlap content below)
    tb_head = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(11.7), Inches(0.65))
    tf_head = tb_head.text_frame; tf_head.word_wrap = True
    p_head = tf_head.paragraphs[0]
    r_head = p_head.add_run()
    r_head.text = main_headline
    r_head.font.size = Pt(25)
    r_head.font.color.rgb = DARK_NAVY
    r_head.font.bold = True
    r_head.font.name = "Segoe UI"

def add_footer(slide):
    """Subtle bottom branding bar"""
    y = H - Inches(0.42)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF6)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.06), W - Inches(1.6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Ishara AI   •   Alibaba Cloud AI Hackathon Pakistan 2026   •   Confidential"
    r.font.size = Pt(9.5)
    r.font.color.rgb = TEXT_MUTED

def card_shape(slide, left, top, w, h, fill=WHITE, border=CARD_BORDER):
    """Draw a clean card background shape"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1)
    return s


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (Clean, high-contrast, beautiful branding)
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)

# Top green accent
top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = GREEN; top_bar.line.fill.background()

# Elegant center container card for title content
card_shape(s1, Inches(2.2), Inches(0.6), Inches(8.933), Inches(6.0), fill=BG_LIGHT, border=GREEN_BORDER)

# Logo Mark centered at top (pure hand mark, transparent, no text collision)
if os.path.exists(LOGO_MARK):
    s1.shapes.add_picture(LOGO_MARK, Inches(5.9), Inches(0.95), height=Inches(1.65))

# "Ishara AI" brand name in sharp typography
tb_name = s1.shapes.add_textbox(Inches(2.6), Inches(2.75), Inches(8.133), Inches(0.6))
tf_name = tb_name.text_frame; tf_name.word_wrap = True
p_n = tf_name.paragraphs[0]; p_n.alignment = PP_ALIGN.CENTER
r_n = p_n.add_run()
r_n.text = "Ishara AI"
r_n.font.size = Pt(36); r_n.font.bold = True; r_n.font.color.rgb = DARK_NAVY; r_n.font.name = "Segoe UI"

# Tagline
tb_tag = s1.shapes.add_textbox(Inches(2.6), Inches(3.45), Inches(8.133), Inches(0.5))
tf_tag = tb_tag.text_frame; tf_tag.word_wrap = True
p_t = tf_tag.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
r_t = p_t.add_run()
r_t.text = "The First Minute When No Interpreter Is Present"
r_t.font.size = Pt(20); r_t.font.bold = True; r_t.font.color.rgb = GREEN; r_t.font.name = "Segoe UI"

# Subtitle
tb_sub = s1.shapes.add_textbox(Inches(2.6), Inches(4.05), Inches(8.133), Inches(0.45))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_s = tf_sub.paragraphs[0]; p_s.alignment = PP_ALIGN.CENTER
r_s = p_s.add_run()
r_s.text = "AI-powered Pakistan Sign Language recognition for emergency healthcare triage"
r_s.font.size = Pt(13.5); r_s.font.color.rgb = TEXT_BODY; r_s.font.name = "Segoe UI"

# Green decorative rule
rule = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(4.7), Inches(2.133), Inches(0.04))
rule.fill.solid(); rule.fill.fore_color.rgb = GREEN; rule.line.fill.background()

# Badges line
tb_b = s1.shapes.add_textbox(Inches(2.6), Inches(4.9), Inches(8.133), Inches(0.4))
tf_b = tb_b.text_frame; tf_b.word_wrap = True
p_b = tf_b.paragraphs[0]; p_b.alignment = PP_ALIGN.CENTER
r_b = p_b.add_run()
r_b.text = "Alibaba Cloud AI Hackathon Pakistan 2026   |   Team Submission"
r_b.font.size = Pt(12); r_b.font.color.rgb = TEXT_MUTED; r_b.font.bold = True; r_b.font.name = "Segoe UI"

# Status badge
tb_st = s1.shapes.add_textbox(Inches(2.6), Inches(5.35), Inches(8.133), Inches(0.4))
tf_st = tb_st.text_frame; tf_st.word_wrap = True
p_st = tf_st.paragraphs[0]; p_st.alignment = PP_ALIGN.CENTER
r_st = p_st.add_run()
r_st.text = "100% Offline   •   No Internet Required   •   Zero Hallucination"
r_st.font.size = Pt(11); r_st.font.color.rgb = GREEN; r_st.font.bold = True

add_footer(s1)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM (No text overlap, clean single-frame cards)
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, BG_LIGHT)
add_header(s2, "01", "THE PROBLEM", "Who is left behind when no interpreter is present?")
add_footer(s2)

stats_data = [
    ("1.1M+", "Deaf & hard-of-hearing", "citizens in Pakistan with virtually zero accessible hospital communication tools."),
    ("< 50", "Certified PSL Interpreters", "nationwide. Rarely available 24/7, and virtually absent in emergency departments."),
    ("0", "ER Communication Aids", "available today in public hospitals. Between patient and doctor, there is only silence."),
]

# 3 Stat cards placed from Y = 1.7 to Y = 4.45 (Height = 2.75)
for i, (num_str, label_str, desc_str) in enumerate(stats_data):
    cx = Inches(0.8 + i * 4.05)
    card_shape(s2, cx, Inches(1.7), Inches(3.6), Inches(2.75), fill=WHITE, border=CARD_BORDER)

    # Use a SINGLE text frame so elements flow naturally and NEVER collide!
    tb = s2.shapes.add_textbox(cx + Inches(0.2), Inches(1.85), Inches(3.2), Inches(2.45))
    tf = tb.text_frame; tf.word_wrap = True

    # Big Number
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.space_after = Pt(6)
    r1 = p1.add_run()
    r1.text = num_str; r1.font.size = Pt(36); r1.font.bold = True
    r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"

    # Category label
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(6)
    r2 = p2.add_run()
    r2.text = label_str; r2.font.size = Pt(13); r2.font.bold = True
    r2.font.color.rgb = DARK_NAVY; r2.font.name = "Segoe UI"

    # Description
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc_str; r3.font.size = Pt(11)
    r3.font.color.rgb = TEXT_MUTED; r3.font.name = "Segoe UI"

# Scenario card at bottom from Y = 4.7 to Y = 6.45 (Height = 1.75)
card_shape(s2, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.75), fill=AMBER_BG, border=AMBER_BORDER)

tb_sc = s2.shapes.add_textbox(Inches(1.15), Inches(4.85), Inches(11.0), Inches(1.45))
tf_sc = tb_sc.text_frame; tf_sc.word_wrap = True

p_sc1 = tf_sc.paragraphs[0]; p_sc1.space_after = Pt(4)
r_sc1 = p_sc1.add_run()
r_sc1.text = "THE REALITY AT 2:00 AM IN AN EMERGENCY ROOM"
r_sc1.font.size = Pt(12); r_sc1.font.bold = True; r_sc1.font.color.rgb = AMBER_DARK

p_sc2 = tf_sc.add_paragraph()
r_sc2 = p_sc2.add_run()
r_sc2.text = "A Deaf patient arrives at Civil Hospital Karachi experiencing severe pain. There is no interpreter on duty. " \
             "The emergency physician does not know Pakistan Sign Language. Today, what happens? Gesturing, confusion, " \
             "guessing, or the patient leaving without treatment. This is not a translation problem—it is an emergency triage crisis."
r_sc2.font.size = Pt(12.5); r_sc2.font.color.rgb = DARK_NAVY; r_sc2.font.name = "Segoe UI"


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION (Screenshot perfectly framed, no overlap)
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, BG_LIGHT)
add_header(s3, "02", "OUR SOLUTION", "60 seconds to the first verified medical message")
add_footer(s3)

# App screenshot container: X=0.8, Y=1.7, W=6.7, H=4.9
card_shape(s3, Inches(0.8), Inches(1.7), Inches(6.7), Inches(4.9), fill=WHITE, border=GREEN_BORDER)
if os.path.exists(SS_APP):
    # Scale screenshot to fit neatly inside card with balanced margins
    s3.shapes.add_picture(SS_APP, Inches(0.95), Inches(1.85), width=Inches(6.4))

# Caption below screenshot inside card
tb_cap = s3.shapes.add_textbox(Inches(0.95), Inches(5.95), Inches(6.4), Inches(0.5))
p_cap = tb_cap.text_frame.paragraphs[0]; p_cap.alignment = PP_ALIGN.CENTER
r_cap = p_cap.add_run()
r_cap.text = "Live Ishara AI interface: Patient signs in front of webcam -> Reviewed Urdu displayed"
r_cap.font.size = Pt(10); r_cap.font.color.rgb = TEXT_MUTED; r_cap.font.bold = True

# Right Side: How It Works column (X=7.8, W=4.7, Y=1.7)
card_shape(s3, Inches(7.8), Inches(1.7), Inches(4.7), Inches(4.9), fill=WHITE, border=CARD_BORDER)

# Header in card
tb_rh = s3.shapes.add_textbox(Inches(8.1), Inches(1.85), Inches(4.1), Inches(0.4))
p_rh = tb_rh.text_frame.paragraphs[0]
r_rh = p_rh.add_run(); r_rh.text = "HOW IT WORKS IN 5 STEPS"; r_rh.font.size = Pt(12)
r_rh.font.bold = True; r_rh.font.color.rgb = GREEN

flow_steps = [
    ("1", "Patient signs before camera", "No login, email, or password required. Zero friction in pain."),
    ("2", "MediaPipe extracts landmarks", "Tracks 33 body pose points and 42 hand landmarks in real time."),
    ("3", "DTW matches motion shape", "Dynamic Time Warping matches motion trajectory against verified PSL."),
    ("4", "Reviewed Urdu text displayed", "Pre-written, clinically audited medical phrase. No hallucination."),
    ("5", "Patient confirms, audio plays", "Patient taps speak. Kokoro TTS speaks Urdu to the doctor."),
]

for idx, (step_n, step_t, step_d) in enumerate(flow_steps):
    step_y = Inches(2.35 + idx * 0.8)

    # Number circle badge
    circ = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.1), step_y + Inches(0.04), Inches(0.36), Inches(0.36))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN; circ.line.fill.background()
    p_c = circ.text_frame.paragraphs[0]; p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run(); r_c.text = step_n; r_c.font.size = Pt(11); r_c.font.color.rgb = WHITE; r_c.font.bold = True

    # Step details in single text box
    tb_s = s3.shapes.add_textbox(Inches(8.6), step_y, Inches(3.7), Inches(0.75))
    tf_s = tb_s.text_frame; tf_s.word_wrap = True
    p_st = tf_s.paragraphs[0]; p_st.space_after = Pt(2)
    r_st = p_st.add_run(); r_st.text = step_t; r_st.font.size = Pt(11.5); r_st.font.bold = True
    r_st.font.color.rgb = DARK_NAVY

    p_sd = tf_s.add_paragraph()
    r_sd = p_sd.add_run(); r_sd.text = step_d; r_sd.font.size = Pt(10); r_sd.font.color.rgb = TEXT_MUTED


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — NEED & IMPACT (Before vs After comparison)
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, BG_LIGHT)
add_header(s4, "03", "NEED & IMPACT", "Not replacing interpreters—bridging the critical gap until one arrives")
add_footer(s4)

# WITHOUT CARD (Red)
card_shape(s4, Inches(0.8), Inches(1.7), Inches(5.65), Inches(4.3), fill=RED_BG, border=RED_BORDER)
tb_w = s4.shapes.add_textbox(Inches(1.15), Inches(1.9), Inches(5.0), Inches(3.8))
tf_w = tb_w.text_frame; tf_w.word_wrap = True

p_wh = tf_w.paragraphs[0]; p_wh.space_after = Pt(12)
r_wh = p_wh.add_run(); r_wh.text = "WITHOUT ISHARA AI  (CURRENT REALITY)"
r_wh.font.size = Pt(13); r_wh.font.bold = True; r_wh.font.color.rgb = RED_TEXT

without_points = [
    ("✗  Unstructured Gesturing:", "Misunderstandings happen immediately. Patient points; doctor guesses."),
    ("✗  Doctor Cannot Ask Questions:", "No two-way interaction. Doctor cannot check pain location or duration."),
    ("✗  Diagnostic & Triage Delay:", "Waiting for an interpreter takes hours or days. Critical windows are lost."),
    ("✗  High Rate of Treatment Abandonment:", "Deaf patients frequently leave ER without receiving care due to frustration."),
]

for title_pt, desc_pt in without_points:
    p_t = tf_w.add_paragraph(); p_t.space_after = Pt(1)
    r_t = p_t.add_run(); r_t.text = title_pt; r_t.font.size = Pt(11.5); r_t.font.bold = True
    r_t.font.color.rgb = RED_TEXT
    p_d = tf_w.add_paragraph(); p_d.space_after = Pt(8)
    r_d = p_d.add_run(); r_d.text = desc_pt; r_d.font.size = Pt(10.5); r_d.font.color.rgb = TEXT_BODY

# WITH CARD (Green)
card_shape(s4, Inches(6.88), Inches(1.7), Inches(5.65), Inches(4.3), fill=GREEN_TINT, border=GREEN_BORDER)
tb_g = s4.shapes.add_textbox(Inches(7.2), Inches(1.9), Inches(5.0), Inches(3.8))
tf_g = tb_g.text_frame; tf_g.word_wrap = True

p_gh = tf_g.paragraphs[0]; p_gh.space_after = Pt(12)
r_gh = p_gh.add_run(); r_gh.text = "WITH ISHARA AI  (IMMEDIATE TRIAGE)"
r_gh.font.size = Pt(13); r_gh.font.bold = True; r_gh.font.color.rgb = GREEN_DARK

with_points = [
    ("✓  First Minute Triage:", "Patient signs -> reviewed Urdu output displayed and spoken in under 60 seconds."),
    ("✓  Two-Way Doctor Responses:", "Doctor uses curated PSL response library to send structured video queries."),
    ("✓  Zero Guessing Safety Gate:", "DTW distance gate refuses classification when unsure: 'Pata nahi, dobara karein'."),
    ("✓  100% Offline & Private:", "Runs entirely on a standard hospital laptop. No patient data stored or uploaded."),
]

for title_pt, desc_pt in with_points:
    p_t = tf_g.add_paragraph(); p_t.space_after = Pt(1)
    r_t = p_t.add_run(); r_t.text = title_pt; r_t.font.size = Pt(11.5); r_t.font.bold = True
    r_t.font.color.rgb = GREEN_DARK
    p_d = tf_g.add_paragraph(); p_d.space_after = Pt(8)
    r_d = p_d.add_run(); r_d.text = desc_pt; r_d.font.size = Pt(10.5); r_d.font.color.rgb = TEXT_BODY

# Impact Quote Banner at Y=6.15 to Y=6.75
card_shape(s4, Inches(0.8), Inches(6.15), Inches(11.73), Inches(0.6), fill=GREEN, border=GREEN)
tb_q = s4.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(11.33), Inches(0.5))
p_q = tb_q.text_frame.paragraphs[0]; p_q.alignment = PP_ALIGN.CENTER
r_q = p_q.add_run()
r_q.text = '"One sentence—‘میری آنکھ میں درد ہے’—delivered in the first 60 seconds changes a clinical triage outcome."'
r_q.font.size = Pt(13); r_q.font.color.rgb = WHITE; r_q.font.bold = True; r_q.font.name = "Segoe UI"


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — INNOVATION & TECHNOLOGY (Architecture + Key Innovations)
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, BG_LIGHT)
add_header(s5, "04", "INNOVATION & TECHNOLOGY", "The system that refuses to lie: clinical safety by design")
add_footer(s5)

# Pipeline cards across top (Y = 1.7 to 3.4)
pipeline_steps = [
    ("📹 Camera", "Webcam Feed", "30 FPS real-time capture via browser"),
    ("🦴 MediaPipe", "Pose & Hands", "33 body + 42 hand 3D landmarks"),
    ("📐 DTW Matcher", "Time Warping", "Matches motion shape regardless of speed"),
    ("🛡️ Safety Gate", "Unknown Reject", "Rejects match if distance > tau threshold"),
    ("📖 Phrase Map", "Curated Urdu", "Look-up table: zero generative AI risk"),
    ("🔊 Kokoro TTS", "Audio Voice", "High-clarity local speech to medical staff"),
]

for i, (p_title, p_sub, p_detail) in enumerate(pipeline_steps):
    px = Inches(0.8 + i * 1.98)
    card_shape(s5, px, Inches(1.7), Inches(1.85), Inches(1.65), fill=WHITE, border=CARD_BORDER)

    tb_p = s5.shapes.add_textbox(px + Inches(0.08), Inches(1.78), Inches(1.7), Inches(1.5))
    tf_p = tb_p.text_frame; tf_p.word_wrap = True

    p_pt = tf_p.paragraphs[0]; p_pt.alignment = PP_ALIGN.CENTER; p_pt.space_after = Pt(2)
    r_pt = p_pt.add_run(); r_pt.text = p_title; r_pt.font.size = Pt(11.5); r_pt.font.bold = True
    r_pt.font.color.rgb = GREEN_DARK

    p_ps = tf_p.add_paragraph(); p_ps.alignment = PP_ALIGN.CENTER; p_ps.space_after = Pt(3)
    r_ps = p_ps.add_run(); r_ps.text = p_sub; r_ps.font.size = Pt(10.5); r_ps.font.bold = True
    r_ps.font.color.rgb = DARK_NAVY

    p_pd = tf_p.add_paragraph(); p_pd.alignment = PP_ALIGN.CENTER
    r_pd = p_pd.add_run(); r_pd.text = p_detail; r_pd.font.size = Pt(9.5); r_pd.font.color.rgb = TEXT_MUTED

    # Arrow to next step
    if i < len(pipeline_steps) - 1:
        tb_arr = s5.shapes.add_textbox(px + Inches(1.76), Inches(2.25), Inches(0.3), Inches(0.3))
        p_arr = tb_arr.text_frame.paragraphs[0]; p_arr.alignment = PP_ALIGN.CENTER
        r_arr = p_arr.add_run(); r_arr.text = "›"; r_arr.font.size = Pt(20); r_arr.font.color.rgb = GREEN; r_arr.font.bold = True

# 4 Key Innovations Grid (Y = 3.6 to Y = 6.45)
innovations = [
    ("🛡️ Zero Hallucination Guarantee",
     "No Large Language Model generates medical responses. Signs map deterministically to clinically reviewed Urdu phrases. "
     "In healthcare, an invented sentence can be fatal. Our architecture makes hallucination mathematically impossible."),

    ("✋ Honest Unknown Gate ('Pata Nahi')",
     "Standard neural classifiers force a prediction (softmax sum = 1). DTW measures absolute Euclidean distance. "
     "If the closest sign is too far (tau > 0.25), the system responds 'Pata nahi, dobara karein' rather than guessing."),

    ("⚡ 100% Local & Edge Feasible",
     "The entire pipeline—MediaPipe Holistic, DTW matcher, and Kokoro TTS—runs locally on a standard CPU. "
     "No internet connection is needed in the hospital. The entire live demo functions with Wi-Fi disabled."),

    ("Built for PSL, Not Generic ASL",
     "Pakistan Sign Language is grammatically and lexically independent from ASL or BSL. Web models trained on ASL fail in Pakistan. "
     "We calibrated against verified PSL references, accounting for distinct Pakistani handshapes and orientations."),
]

for idx, (inv_title, inv_desc) in enumerate(innovations):
    row = idx // 2
    col = idx % 2
    ix = Inches(0.8 + col * 5.95)
    iy = Inches(3.6 + row * 1.45)

    card_shape(s5, ix, iy, Inches(5.78), Inches(1.35), fill=WHITE, border=CARD_BORDER)

    tb_inv = s5.shapes.add_textbox(ix + Inches(0.2), iy + Inches(0.12), Inches(5.38), Inches(1.15))
    tf_inv = tb_inv.text_frame; tf_inv.word_wrap = True

    p_it = tf_inv.paragraphs[0]; p_it.space_after = Pt(3)
    r_it = p_it.add_run(); r_it.text = inv_title; r_it.font.size = Pt(12); r_it.font.bold = True
    r_it.font.color.rgb = GREEN_DARK

    p_id = tf_inv.add_paragraph()
    r_id = p_id.add_run(); r_id.text = inv_desc; r_id.font.size = Pt(10.5); r_id.font.color.rgb = TEXT_BODY


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — FEASIBILITY & WHAT WE BUILT (No text overflow!)
# ═══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, BG_LIGHT)
add_header(s6, "05", "FEASIBILITY & WHAT WE BUILT", "Working software with measured metrics, not a conceptual mockup")
add_footer(s6)

# Left Card: Measured Accuracy (X=0.8, Y=1.7, W=5.7, H=3.2)
card_shape(s6, Inches(0.8), Inches(1.7), Inches(5.7), Inches(3.2), fill=GREEN_TINT, border=GREEN_BORDER)
tb_m = s6.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(5.1), Inches(2.9))
tf_m = tb_m.text_frame; tf_m.word_wrap = True

p_mh = tf_m.paragraphs[0]; p_mh.space_after = Pt(8)
r_mh = p_mh.add_run(); r_mh.text = "VERIFIED METRICS (WITH DENOMINATORS)"
r_mh.font.size = Pt(12.5); r_mh.font.bold = True; r_mh.font.color.rgb = GREEN_DARK

metrics_list = [
    ("Leave-One-Out Validation:", "24 / 24 correct (100% on 4 core signs across 24 reference clips)"),
    ("Live Interactive Session:", "14 correct, 0 incorrect, 3 unknown (100% clean safety gate)"),
    ("Margin of Separation:", "Same-sign worst: 0.161 vs Cross-sign best: 0.191 (clear decision gap)"),
    ("False Positive Rate:", "< 2% target maintained by design; rejects rather than misclassifies"),
    ("Honest Test Context:", "Measured on single deliberate signer (P01); ready for multi-signer Phase 2"),
]

for label_m, val_m in metrics_list:
    p_m = tf_m.add_paragraph(); p_m.space_after = Pt(4)
    r_ml = p_m.add_run(); r_ml.text = label_m + " "; r_ml.font.size = Pt(10.5); r_ml.font.bold = True; r_ml.font.color.rgb = DARK_NAVY
    r_mv = p_m.add_run(); r_mv.text = val_m; r_mv.font.size = Pt(10); r_mv.font.color.rgb = TEXT_BODY

# Right Card: What Is Running Today (X=6.83, Y=1.7, W=5.7, H=3.2)
card_shape(s6, Inches(6.83), Inches(1.7), Inches(5.7), Inches(3.2), fill=WHITE, border=CARD_BORDER)
tb_b = s6.shapes.add_textbox(Inches(7.1), Inches(1.85), Inches(5.2), Inches(2.9))
tf_b = tb_b.text_frame; tf_b.word_wrap = True

p_bh = tf_b.paragraphs[0]; p_bh.space_after = Pt(8)
r_bh = p_bh.add_run(); r_bh.text = "FULL-STACK PRODUCTION SYSTEM RUNNING"
r_bh.font.size = Pt(12.5); r_bh.font.bold = True; r_bh.font.color.rgb = DARK_NAVY

built_points = [
    "✓ Full-Stack Web App: Next.js 14 frontend + FastAPI Python backend + PostgreSQL",
    "✓ Real-time webcam processing: 30 FPS landmark extraction in client browser",
    "✓ Dynamic Time Warping engine: Subsequence matching against calibrated references",
    "✓ Doctor Mode with PSL Video Library: Spoken/button questions trigger verified PSL video",
    "✓ High-clarity Urdu Voice: Kokoro phonetic Hindi-Urdu pipeline checked by ear",
    "✓ Resilient Fallback Chain: Database -> JSON snapshot; Cloud speech -> Local speech",
]

for pt in built_points:
    p_bp = tf_b.add_paragraph(); p_bp.space_after = Pt(4)
    r_bp = p_bp.add_run(); r_bp.text = pt; r_bp.font.size = Pt(10); r_bp.font.color.rgb = TEXT_BODY

# Bottom Card: Honest Limitations (X=0.8, Y=5.1, W=11.73, H=1.6)
card_shape(s6, Inches(0.8), Inches(5.1), Inches(11.73), Inches(1.6), fill=AMBER_BG, border=AMBER_BORDER)

tb_lh = s6.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(11.1), Inches(0.3))
p_lh = tb_lh.text_frame.paragraphs[0]
r_lh = p_lh.add_run(); r_lh.text = "TRANSPARENT LIMITATIONS — WE STATE THESE FIRST BEFORE JUDGES ASK"
r_lh.font.size = Pt(11.5); r_lh.font.bold = True; r_lh.font.color.rgb = AMBER_DARK

lims = [
    ("Fluent Continuous Signing:", "Currently requires ~320ms stillness between signs. Sliding-window spotting is Phase 2 roadmap."),
    ("Full Conversation Syntax:", "Global open research challenge. We focus on urgent triage statements, not full linguistic translation."),
    ("Not an Interpreter Replacement:", "Complex clinical examinations still require certified interpreters. We unlock the first critical minutes."),
]

for i, (lim_t, lim_d) in enumerate(lims):
    lx = Inches(1.1 + i * 3.8)
    tb_l = s6.shapes.add_textbox(lx, Inches(5.55), Inches(3.6), Inches(1.0))
    tf_l = tb_l.text_frame; tf_l.word_wrap = True

    p_lt = tf_l.paragraphs[0]; p_lt.space_after = Pt(2)
    r_lt = p_lt.add_run(); r_lt.text = "▸ " + lim_t; r_lt.font.size = Pt(10.5); r_lt.font.bold = True; r_lt.font.color.rgb = DARK_NAVY

    p_ld = tf_l.add_paragraph()
    r_ld = p_ld.add_run(); r_ld.text = lim_d; r_ld.font.size = Pt(9.5); r_ld.font.color.rgb = TEXT_MUTED


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — CLOSING & IMPACT (Clean, professional finale)
# ═══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, WHITE)

# Top green accent
t_bar7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
t_bar7.fill.solid(); t_bar7.fill.fore_color.rgb = GREEN; t_bar7.line.fill.background()

# Center Card
card_shape(s7, Inches(2.2), Inches(0.6), Inches(8.933), Inches(6.0), fill=BG_LIGHT, border=CARD_BORDER)

# Logo mark centered at top
if os.path.exists(LOGO_MARK):
    s7.shapes.add_picture(LOGO_MARK, Inches(5.9), Inches(0.95), height=Inches(1.65))

# Brand Name
tb_n7 = s7.shapes.add_textbox(Inches(2.6), Inches(2.75), Inches(8.133), Inches(0.55))
tf_n7 = tb_n7.text_frame; tf_n7.word_wrap = True
p_n7 = tf_n7.paragraphs[0]; p_n7.alignment = PP_ALIGN.CENTER
r_n7 = p_n7.add_run()
r_n7.text = "Ishara AI"
r_n7.font.size = Pt(36); r_n7.font.bold = True; r_n7.font.color.rgb = DARK_NAVY; r_n7.font.name = "Segoe UI"

# Tagline
tb_t7 = s7.shapes.add_textbox(Inches(2.6), Inches(3.35), Inches(8.133), Inches(0.5))
tf_t7 = tb_t7.text_frame; tf_t7.word_wrap = True
p_t7 = tf_t7.paragraphs[0]; p_t7.alignment = PP_ALIGN.CENTER
r_t7 = p_t7.add_run()
r_t7.text = "The first minute. Not the last word."
r_t7.font.size = Pt(20); r_t7.font.bold = True; r_t7.font.color.rgb = GREEN; r_t7.font.name = "Segoe UI"

# Subtitle
tb_s7 = s7.shapes.add_textbox(Inches(2.6), Inches(3.95), Inches(8.133), Inches(0.45))
tf_s7 = tb_s7.text_frame; tf_s7.word_wrap = True
p_s7 = tf_s7.paragraphs[0]; p_s7.alignment = PP_ALIGN.CENTER
r_s7 = p_s7.add_run()
r_s7.text = "Empowering Pakistani healthcare with honest, offline-first assistive AI"
r_s7.font.size = Pt(13.5); r_s7.font.color.rgb = TEXT_BODY; r_s7.font.name = "Segoe UI"

# GitHub & Repo Box
repo_box = card_shape(s7, Inches(3.0), Inches(4.55), Inches(7.333), Inches(0.65), fill=WHITE, border=GREEN_BORDER)
tb_repo = s7.shapes.add_textbox(Inches(3.1), Inches(4.65), Inches(7.133), Inches(0.45))
p_repo = tb_repo.text_frame.paragraphs[0]; p_repo.alignment = PP_ALIGN.CENTER
r_repo = p_repo.add_run()
r_repo.text = "GitHub: Obaid35/Ishara-AI---ALIBABA-CLOUD-AI-HACKATHON-PAKISTAN-2026"
r_repo.font.size = Pt(11.5); r_repo.font.bold = True; r_repo.font.color.rgb = DARK_NAVY

# Closing prompt
tb_demo = s7.shapes.add_textbox(Inches(2.6), Inches(5.35), Inches(8.133), Inches(0.5))
p_demo = tb_demo.text_frame.paragraphs[0]; p_demo.alignment = PP_ALIGN.CENTER
r_demo = p_demo.add_run()
r_demo.text = "Thank You   •   Live System Demo Ready for Judges"
r_demo.font.size = Pt(15); r_demo.font.bold = True; r_demo.font.color.rgb = GREEN_DARK

add_footer(s7)

# ── Save Presentation ──
prs.save(OUTPUT)
print(f"SUCCESS: Presentation saved -> {OUTPUT}")
